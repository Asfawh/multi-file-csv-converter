import os
import csv
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from docx import Document
from subprocess import run, CalledProcessError

# Directory paths
pdf_directory = 'Mezmure/PDF'
ppt_directory = 'Mezmure/PPT'
doc_directory = 'Mezmure/WORD'
output_directory = 'Mezmure/CSV_Output0'
error_log_file = os.path.join(output_directory, 'error_log.txt')

# Ensure output directory exists
if not os.path.exists(output_directory):
    os.makedirs(output_directory)

# Error logging setup
def log_error(message):
    with open(error_log_file, 'a') as error_log:
        error_log.write(message + '\n')

# Function to add standard columns to each row
def add_standard_columns(file_name, page_number, text, genre="Unknown Genre", zemari_name="Unknown Zemari"):
    mezmur_name = os.path.splitext(os.path.basename(file_name))[0]  # Extract file name without extension
    file_name_only = os.path.basename(file_name)  # Only file name without path
    return [mezmur_name, page_number, zemari_name, file_name_only, genre, text]

# PDF to CSV conversion
def process_pdf_file(pdf_file_path, output_directory):
    try:
        pages = convert_from_path(pdf_file_path)
        data = []
        for page_number, page in enumerate(pages, start=1):
            text = pytesseract.image_to_string(page, lang='amh')
            data.append(add_standard_columns(pdf_file_path, page_number, text))

        csv_file_name = os.path.splitext(os.path.basename(pdf_file_path))[0] + '.csv'
        csv_file_path = os.path.join(output_directory, csv_file_name)

        # Write the data with headers
        with open(csv_file_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['Mezmur Name', 'Page/Verse', 'Zemari Name', 'File Name', 'Genre', 'Verses'])  # Add headers
            writer.writerows(data)
        print(f"Converted {pdf_file_path} to {csv_file_name}")
    except Exception as e:
        log_error(f"Error processing PDF file {pdf_file_path}: {str(e)}")

# PPT to CSV conversion
def process_ppt_file(ppt_file_path, output_directory):
    try:
        prs = Presentation(ppt_file_path)
        data = []
        for page_number, slide in enumerate(prs.slides, start=1):
            slide_text = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, 'text')])
            data.append(add_standard_columns(ppt_file_path, page_number, slide_text))

        csv_file_name = os.path.splitext(os.path.basename(ppt_file_path))[0] + '.csv'
        csv_file_path = os.path.join(output_directory, csv_file_name)

        with open(csv_file_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['Mezmur Name', 'Page/Verse', 'Zemari Name', 'File Name', 'Genre', 'Verses'])
            writer.writerows(data)
        print(f"Converted {ppt_file_path} to {csv_file_name}")
    except Exception as e:
        log_error(f"Error processing PPT file {ppt_file_path}: {str(e)}")

# DOCX to CSV conversion
def process_docx_file(docx_file_path, output_directory):
    try:
        doc = Document(docx_file_path)
        data = []
        page_number = 1
        paragraph_group = []  # Will hold text until a page-like boundary is encountered (e.g., blank line)

        # Loop through each paragraph
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():  # If paragraph contains text, add to group
                paragraph_group.append(paragraph.text)
            else:  # When encountering a blank line, treat it as a new page or section
                if paragraph_group:
                    # Join the paragraphs in the group, treat as one row for this "page"
                    combined_text = '\n'.join(paragraph_group)
                    data.append(add_standard_columns(docx_file_path, page_number, combined_text))
                    paragraph_group = []  # Reset group for the next page
                    page_number += 1
        
        # Catch any remaining text if the last page doesn't end in a blank line
        if paragraph_group:
            combined_text = '\n'.join(paragraph_group)
            data.append(add_standard_columns(docx_file_path, page_number, combined_text))

        # Write CSV file
        csv_file_name = os.path.splitext(os.path.basename(docx_file_path))[0] + '.csv'
        csv_file_path = os.path.join(output_directory, csv_file_name)

        with open(csv_file_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(['Mezmur Name', 'Page/Verse', 'Zemari Name', 'File Name', 'Genre', 'Verses'])  # Add headers
            writer.writerows(data)
        
        print(f"Converted {docx_file_path} to {csv_file_name}")

    except Exception as e:
        log_error(f"Error processing DOCX file {docx_file_path}: {str(e)}")


# Convert .doc to .docx
def convert_doc_to_docx(doc_file_path):
    try:
        docx_file_path = os.path.splitext(doc_file_path)[0] + '.docx'
        run(['soffice', '--headless', '--convert-to', 'docx', doc_file_path, '--outdir', os.path.dirname(doc_file_path)], check=True)
        print(f"Converted {doc_file_path} to {docx_file_path}")
        return docx_file_path
    except CalledProcessError as e:
        log_error(f"Error converting DOC to DOCX {doc_file_path}: {str(e)}")
        return None

# Main processing function to handle all file types
def process_files():
    # Process PDFs
    for pdf_file in os.listdir(pdf_directory):
        if pdf_file.endswith('.pdf'):
            pdf_file_path = os.path.join(pdf_directory, pdf_file)
            process_pdf_file(pdf_file_path, output_directory)

    # Process PPT files
    for ppt_file in os.listdir(ppt_directory):
        if ppt_file.endswith('.pptx'):
            ppt_file_path = os.path.join(ppt_directory, ppt_file)
            process_ppt_file(ppt_file_path, output_directory)

    # Process DOC and DOCX files
    for doc_file in os.listdir(doc_directory):
        doc_file_path = os.path.join(doc_directory, doc_file)
        if doc_file.endswith('.docx'):
            process_docx_file(doc_file_path, output_directory)
        elif doc_file.endswith('.doc'):
            docx_file_path = convert_doc_to_docx(doc_file_path)
            if docx_file_path:
                process_docx_file(docx_file_path, output_directory)

# Run the process
process_files()

print("All files have been processed.")
